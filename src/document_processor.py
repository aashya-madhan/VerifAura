"""
document_processor.py — Convert document files to images for ELA analysis.

Supported formats:
  PDF  — via PyMuPDF (fitz)
  DOCX — via python-docx + HTML rasterization (falls back to embedded images)
  PPTX — via python-pptx (extracts slide images or renders via PIL)

Each document is converted into a list of PIL Images (one per page/slide),
which are then saved to disk as temporary PNGs so the existing ELA pipeline
can process them unchanged.
"""

from __future__ import annotations

import io
import os
import tempfile
from pathlib import Path
from typing import List, Tuple

from PIL import Image

# ──────────────────────────────────────────────────────────────
# PDF support (PyMuPDF)
# ──────────────────────────────────────────────────────────────
def _pdf_to_images(path: str, dpi: int = 150) -> List[Image.Image]:
    """Rasterize every page of a PDF to an RGB PIL Image."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError(
            "PyMuPDF is required for PDF support. "
            "Install it with: pip install pymupdf"
        )

    doc = fitz.open(path)
    images: List[Image.Image] = []
    zoom = dpi / 72.0  # 72 is PDF's default unit
    mat = fitz.Matrix(zoom, zoom)

    for page in doc:
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        images.append(img)

    doc.close()
    return images


# ──────────────────────────────────────────────────────────────
# DOCX support (python-docx)
# ──────────────────────────────────────────────────────────────
def _docx_to_images(path: str) -> List[Image.Image]:
    """
    Extract embedded images from a DOCX file.
    Each image in the document becomes one analysis target.
    If there are no embedded images, we render the first page as a placeholder message.
    """
    try:
        import docx  # python-docx
    except ImportError:
        raise ImportError(
            "python-docx is required for DOCX support. "
            "Install it with: pip install python-docx"
        )

    doc = docx.Document(path)
    images: List[Image.Image] = []

    # Walk all inline images embedded in the document
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                img = Image.open(io.BytesIO(img_data)).convert("RGB")
                images.append(img)
            except Exception:
                pass  # skip unreadable embedded images

    # If the DOCX has no images, try to create a visual representation
    # by rendering a page-like canvas with the document text
    if not images:
        images = [_render_text_page(doc)]

    return images


def _render_text_page(doc, width: int = 794, height: int = 1123) -> Image.Image:
    """
    Render a simple white-background page image with the first ~50 lines
    of text from a document. Used as fallback when DOCX has no embedded images.
    """
    from PIL import ImageDraw, ImageFont

    img = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Try to use a basic font; fall back to default if not available
    try:
        font = ImageFont.truetype("arial.ttf", size=14)
    except Exception:
        font = ImageFont.load_default()

    y = 40
    x = 40
    max_lines = 60
    lines_drawn = 0

    for para in doc.paragraphs:
        if lines_drawn >= max_lines:
            break
        text = para.text.strip()
        if not text:
            y += 8
            continue

        # Wrap long lines
        words = text.split()
        line = ""
        for word in words:
            test = line + (" " if line else "") + word
            if len(test) * 8 > (width - 2 * x):
                draw.text((x, y), line, fill=(20, 20, 20), font=font)
                y += 20
                lines_drawn += 1
                line = word
                if lines_drawn >= max_lines:
                    break
            else:
                line = test
        if line and lines_drawn < max_lines:
            draw.text((x, y), line, fill=(20, 20, 20), font=font)
            y += 20
            lines_drawn += 1

    return img


# ──────────────────────────────────────────────────────────────
# PPTX support (python-pptx)
# ──────────────────────────────────────────────────────────────
def _pptx_to_images(path: str) -> List[Image.Image]:
    """
    Extract images from each slide in a PPTX file.
    Each slide's embedded images are collected; if a slide has none,
    a text-rendered thumbnail is generated.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX support. "
            "Install it with: pip install python-pptx"
        )

    prs = Presentation(path)
    slide_images: List[Image.Image] = []

    for slide_idx, slide in enumerate(prs.slides):
        found_images: List[Image.Image] = []

        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_data = shape.image.blob
                    img = Image.open(io.BytesIO(img_data)).convert("RGB")
                    found_images.append(img)
                except Exception:
                    pass
            # Also check for images inside group shapes
            elif hasattr(shape, "shapes"):
                for sub in shape.shapes:
                    if sub.shape_type == 13:
                        try:
                            img_data = sub.image.blob
                            img = Image.open(io.BytesIO(img_data)).convert("RGB")
                            found_images.append(img)
                        except Exception:
                            pass

        if found_images:
            slide_images.extend(found_images)
        else:
            # Render a slide thumbnail from text content
            slide_images.append(_render_slide_page(slide, slide_idx + 1))

    return slide_images


def _render_slide_page(slide, slide_num: int, width: int = 960, height: int = 540) -> Image.Image:
    """Create a simple visual for a PPTX slide that has no embedded images."""
    from PIL import ImageDraw, ImageFont

    img = Image.new("RGB", (width, height), color=(240, 240, 245))
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("arial.ttf", size=22)
        font_body  = ImageFont.truetype("arial.ttf", size=14)
    except Exception:
        font_title = ImageFont.load_default()
        font_body  = font_title

    draw.text((30, 20), f"Slide {slide_num}", fill=(100, 100, 120), font=font_body)

    y = 60
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()[:200]
            draw.text((30, y), text, fill=(20, 20, 20), font=font_title if y < 100 else font_body)
            y += 30
            if y > height - 40:
                break

    return img


# ──────────────────────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────────────────────

DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}
IMAGE_EXTENSIONS    = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"}


def is_document(filename: str) -> bool:
    """Return True if the file is a document type (not a plain image)."""
    return Path(filename).suffix.lower() in DOCUMENT_EXTENSIONS


def is_image(filename: str) -> bool:
    """Return True if the file is a plain image."""
    return Path(filename).suffix.lower() in IMAGE_EXTENSIONS


def extract_pages(filepath: str, upload_dir: str) -> List[Tuple[str, int]]:
    """
    Convert a document to a list of page images saved on disk.

    Args:
        filepath:   path to the uploaded document or image
        upload_dir: directory to save the extracted page images

    Returns:
        List of (image_path, page_number) tuples.
        For plain images, returns [(filepath, 1)].

    Raises:
        ValueError: for unsupported file types
        ImportError: if a required library is not installed
    """
    ext = Path(filepath).suffix.lower()
    stem = Path(filepath).stem

    if ext in IMAGE_EXTENSIONS:
        # Plain image — return as-is (page 1)
        return [(filepath, 1)]

    # Convert document to PIL images
    if ext == ".pdf":
        pil_images = _pdf_to_images(filepath)
    elif ext in {".docx", ".doc"}:
        pil_images = _docx_to_images(filepath)
    elif ext in {".pptx", ".ppt"}:
        pil_images = _pptx_to_images(filepath)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # Save each page as PNG
    page_paths: List[Tuple[str, int]] = []
    for i, img in enumerate(pil_images, start=1):
        page_filename = f"{stem}_page{i:03d}.png"
        page_path = os.path.join(upload_dir, page_filename)
        img.save(page_path, "PNG")
        page_paths.append((page_path, i))

    return page_paths


def get_page_count(filepath: str) -> int:
    """Quick helper to count pages in a document without extracting images."""
    ext = Path(filepath).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return 1
    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(filepath)
            n = doc.page_count
            doc.close()
            return n
        except Exception:
            return 1
    if ext in {".docx", ".doc"}:
        return 1  # treat whole doc as one unit for quick count
    if ext in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            return len(Presentation(filepath).slides)
        except Exception:
            return 1
    return 1
